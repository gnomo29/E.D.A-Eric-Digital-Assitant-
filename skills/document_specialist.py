"""Especialista de documentos: informe, CV, traducción y presentaciones."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

try:
    from docx import Document
except Exception:
    Document = None  # type: ignore[assignment]

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except Exception:
    Presentation = None  # type: ignore[assignment]


def create_report_docx(title: str, sections: Iterable[str], output_path: str) -> dict[str, str]:
    if Document is None:
        return {"status": "error", "message": "python-docx no disponible"}
    doc = Document()
    doc.add_heading(title, level=1)
    for section in sections:
        doc.add_paragraph(str(section))
    doc.save(output_path)
    return {"status": "ok", "message": output_path}


def create_cv_docx(name: str, profile: str, output_path: str) -> dict[str, str]:
    if Document is None:
        return {"status": "error", "message": "python-docx no disponible"}
    doc = Document()
    doc.add_heading(name, level=1)
    doc.add_paragraph(profile)
    doc.add_heading("Experiencia", level=2)
    doc.add_paragraph("- Añadir experiencia profesional aquí.")
    doc.add_heading("Educación", level=2)
    doc.add_paragraph("- Añadir formación académica aquí.")
    doc.save(output_path)
    return {"status": "ok", "message": output_path}


def extract_pdf_text(pdf_path: str) -> dict[str, str]:
    if PdfReader is None:
        return {"status": "error", "message": "pypdf no disponible"}
    reader = PdfReader(pdf_path)
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    return {"status": "ok", "message": text[:5000]}


def translate_text_offline(text: str, target_language: str = "es") -> dict[str, str]:
    # Placeholder offline-friendly: delegable a modelo local Ollama vía core.ask.
    translated = f"[{target_language}] {text}"
    return {"status": "ok", "message": translated}


def create_presentation(topic: str, slides: int, output_path: str) -> dict[str, str]:
    if Presentation is None:
        return {"status": "error", "message": "python-pptx no disponible"}
    ppt = Presentation()
    total = max(1, min(20, int(slides)))
    for i in range(total):
        layout = ppt.slide_layouts[1] if len(ppt.slide_layouts) > 1 else ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = f"{topic} - Parte {i + 1}"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Punto clave {i + 1} sobre {topic}."
    ppt.save(output_path)
    return {"status": "ok", "message": output_path}

